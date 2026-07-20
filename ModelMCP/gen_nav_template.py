"""
导览栏 PPT 模板生成器
风格：深蓝导览栏，当前章节亮白文字+青色指示条，其余章节暗色文字
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 配色（蓝色科技风，对标参考图）────────────────────────────
NAV_BG   = RGBColor(0x0C, 0x1E, 0x3C)   # 导览栏底色：极深海蓝
NAV_ACT  = RGBColor(0x14, 0x3A, 0x6E)   # 激活标签背景：中深蓝
CYAN     = RGBColor(0x00, 0xC8, 0xF0)   # 青色指示条 / 装饰线
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x5A, 0x88, 0xB8)   # 未激活标签文字：灰蓝
OFFWHT   = RGBColor(0xF0, 0xF5, 0xFB)   # 内容页底色
SLIDE_BG = RGBColor(0xFA, 0xFC, 0xFF)   # 内容白区
DARK     = RGBColor(0x1A, 0x2E, 0x44)   # 正文深色
LGRAY    = RGBColor(0xD8, 0xE6, 0xF4)   # 分割线
SIDEBAR  = RGBColor(0x0E, 0x4D, 0x92)   # 左侧蓝色竖条
LBL_BLUE = RGBColor(0x9A, 0xB0, 0xFF)  # placeholder - unused
GBLU     = RGBColor(0x1A, 0x5F, 0xAA)   # 渐变蓝（章节页）
WHTFADE  = RGBColor(0xE8, 0xF2, 0xFF)   # 浅蓝白（右侧背景）

SECTIONS = ["问题分析", "矩阵计算", "矩阵压缩", "整体模型", "拓展延伸"]
NUMS     = ["01", "02", "03", "04", "05"]
FONT     = "微软雅黑"

W  = prs.slide_width
H  = prs.slide_height
NH = Inches(0.62)          # 导览栏高度
BL = prs.slide_layouts[6]  # 空白版式

# ── 基础绘图函数 ─────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(9, int(l), int(t), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def label(slide, text, l, t, w, h,
          size=Pt(12), bold=False, color=WHITE,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = tf.margin_bottom = 0
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return tb

# ── 顶部导览栏（核心函数）────────────────────────────────────
def nav_bar(slide, active):
    """
    深蓝底色导览栏。
    激活标签：稍亮蓝底 + 白色粗体 + 底部青色指示条。
    未激活标签：暗灰蓝文字，无底色变化，有竖向分隔线。
    """
    # 整体底色
    rect(slide, 0, 0, W, NH, fill=NAV_BG)

    tab_w   = W / len(SECTIONS)
    ind_h   = Inches(0.055)    # 青色指示条高度

    for i, sec in enumerate(SECTIONS):
        x     = tab_w * i
        is_act = (i == active)

        if is_act:
            # 激活标签：稍亮背景
            rect(slide, x, 0, tab_w, NH, fill=NAV_ACT)
            # 白色文字（粗体）
            label(slide, sec, x, 0, tab_w, NH - ind_h,
                  size=Pt(13), bold=True, color=WHITE)
            # 底部青色指示条
            rect(slide, x, NH - ind_h, tab_w, ind_h, fill=CYAN)
        else:
            # 未激活：暗色文字
            label(slide, sec, x, 0, tab_w, NH,
                  size=Pt(12), bold=False, color=MUTED)

        # 标签间竖向分隔线（最后一个不加右侧线）
        if i < len(SECTIONS) - 1:
            sep_x = x + tab_w - Inches(0.008)
            rect(slide, sep_x, Inches(0.1), Inches(0.012),
                 NH - Inches(0.2),
                 fill=RGBColor(0x1E, 0x3A, 0x5C))

    # 导览栏底部青色细线（全宽）
    rect(slide, 0, NH, W, Inches(0.03), fill=CYAN)

# ── 封面底板 ─────────────────────────────────────────────────
def cover_base(slide):
    """封面/结尾：左侧深蓝面板 + 右侧浅蓝白"""
    rect(slide, 0, 0, W, H, fill=WHTFADE)
    rect(slide, 0, 0, W * 0.38, H, fill=NAV_BG)
    rect(slide, W * 0.38, 0, Inches(0.07), H, fill=CYAN)
    # 圆点网格装饰
    for ri in range(6):
        for ci in range(5):
            oval(slide,
                 Inches(0.55) + ci * Inches(0.50),
                 Inches(0.85) + ri * Inches(0.55),
                 Inches(0.08), Inches(0.08),
                 fill=RGBColor(0x1E, 0x5A, 0x9A))

# ════════════════════════════════════════════════════════════
# 1. 封面
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
cover_base(s)

# 左栏信息
label(s, "汇报人：姓名",
      Inches(0.35), H * 0.76, W * 0.35, Inches(0.38),
      size=Pt(11), color=MUTED, align=PP_ALIGN.LEFT)
label(s, "指导老师：教授姓名",
      Inches(0.35), H * 0.81, W * 0.35, Inches(0.38),
      size=Pt(11), color=MUTED, align=PP_ALIGN.LEFT)
label(s, "2026 年 6 月",
      Inches(0.35), H * 0.87, W * 0.35, Inches(0.38),
      size=Pt(11), color=RGBColor(0x40, 0x70, 0xA0), align=PP_ALIGN.LEFT)

# 右栏标题
tx = W * 0.38 + Inches(0.55)
tw = W * 0.60 - Inches(0.5)
label(s, "演示文稿标题",
      tx, H * 0.20, tw, Inches(1.4),
      size=Pt(38), bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(s, tx, H * 0.20 + Inches(1.42), Inches(2.2), Inches(0.06), fill=CYAN)
label(s, "副标题 · 项目 / 课题名称",
      tx, H * 0.20 + Inches(1.60), tw, Inches(0.72),
      size=Pt(18), color=SIDEBAR, align=PP_ALIGN.LEFT)
label(s, "学院 / 单位  ·  指导老师：教授姓名",
      tx, H * 0.20 + Inches(2.42), tw, Inches(0.5),
      size=Pt(13), color=RGBColor(0x70, 0x90, 0xB0), align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════
# 2. 目录页（Overview — 所有章节均为默认态，无激活）
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
rect(s, 0, 0, W, H, fill=OFFWHT)
nav_bar(s, -1)   # -1 = 无激活（传入不在范围的值）

# 目录标题
tt = NH + Inches(0.04)
rect(s, 0, tt, W, Inches(0.9), fill=SLIDE_BG)
rect(s, 0, tt, Inches(0.08), Inches(0.9), fill=CYAN)
label(s, "目  录", Inches(0.22), tt, W - Inches(0.4), Inches(0.9),
      size=Pt(26), bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(s, Inches(0.22), tt + Inches(0.9), W - Inches(0.44), Inches(0.03), fill=LGRAY)
rect(s, Inches(0.22), tt + Inches(0.9), Inches(1.8), Inches(0.05), fill=CYAN)

# 目录列表（两列）
ct = tt + Inches(1.05)
col_h = Inches(0.85)
col_w = (W - Inches(1.5)) / 2

for i, (sec, num) in enumerate(zip(SECTIONS, NUMS)):
    row = i % 3
    col = i // 3
    cx  = Inches(0.4) + col * (col_w + Inches(0.5))
    cy  = ct + row * (col_h + Inches(0.12))

    rect(s, cx, cy, col_w, col_h, fill=NAV_BG)
    rect(s, cx, cy, Inches(0.07), col_h, fill=CYAN)
    label(s, num,  cx + Inches(0.15), cy, Inches(0.65), col_h,
          size=Pt(22), bold=True, color=CYAN)
    label(s, sec,  cx + Inches(0.8), cy, col_w - Inches(0.9), col_h,
          size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════
# 3. 章节页（分隔 + 内容）× N
# ════════════════════════════════════════════════════════════
for idx, (sec, num) in enumerate(zip(SECTIONS, NUMS)):

    # ── 章节分隔页 ─────────────────────────────────────────
    s = prs.slides.add_slide(BL)
    rect(s, 0, 0, W, H, fill=NAV_BG)
    nav_bar(s, idx)

    # 左侧青色竖条
    rect(s, 0, NH + Inches(0.05), Inches(0.10), H - NH, fill=CYAN)

    # 右侧浅色装饰面板
    rect(s, W * 0.55, NH, W * 0.45, H - NH,
         fill=RGBColor(0x0E, 0x28, 0x50))

    # 水印大数字
    label(s, num,
          W * 0.53, NH + Inches(0.1), W * 0.44, H - NH - Inches(0.15),
          size=Pt(130), bold=True,
          color=RGBColor(0x14, 0x38, 0x68),
          align=PP_ALIGN.RIGHT)

    # 章节名
    label(s, sec,
          Inches(0.35), NH + Inches(1.1), W * 0.58, Inches(1.2),
          size=Pt(44), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 青色分割线
    rect(s, Inches(0.35), NH + Inches(2.45), Inches(3.5), Inches(0.05), fill=CYAN)

    label(s, f"Section {idx + 1}  /  {len(SECTIONS)}",
          Inches(0.35), NH + Inches(2.62), W * 0.5, Inches(0.5),
          size=Pt(13), color=MUTED, align=PP_ALIGN.LEFT)

    # ── 内容页 ─────────────────────────────────────────────
    s = prs.slides.add_slide(BL)
    rect(s, 0, 0, W, H, fill=OFFWHT)
    nav_bar(s, idx)

    # 标题带
    tt = NH + Inches(0.04)
    th = Inches(0.88)
    rect(s, 0, tt, W, th, fill=SLIDE_BG)
    rect(s, 0, tt, Inches(0.08), th, fill=CYAN)
    label(s, sec + "  ·  小节标题",
          Inches(0.22), tt, W - Inches(0.4), th,
          size=Pt(22), bold=True, color=DARK, align=PP_ALIGN.LEFT)
    rect(s, Inches(0.22), tt + th, W - Inches(0.44), Inches(0.025), fill=LGRAY)
    rect(s, Inches(0.22), tt + th, Inches(1.8), Inches(0.05), fill=CYAN)

    # 内容区
    ct  = tt + th + Inches(0.12)
    ch  = H - ct - Inches(0.18)
    cw  = (W - Inches(0.54)) / 2 - Inches(0.1)
    cl  = Inches(0.22)
    cr  = cl + cw + Inches(0.2)

    # 左栏：要点文字
    rect(s, cl, ct, cw, ch, fill=SLIDE_BG)
    rect(s, cl, ct, Inches(0.05), ch, fill=LGRAY)
    label(s,
          "• 核心要点一：在此填写说明文字\n\n"
          "• 核心要点二：在此填写说明文字\n\n"
          "• 核心要点三：在此填写说明文字\n\n"
          "• 核心要点四：在此填写说明文字",
          cl + Inches(0.15), ct + Inches(0.15),
          cw - Inches(0.25), ch - Inches(0.3),
          size=Pt(14), color=DARK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # 右栏：图表占位
    rect(s, cr, ct, cw, ch, fill=RGBColor(0xE0, 0xEC, 0xF8))
    rect(s, cr, ct, cw, Inches(0.05), fill=CYAN)
    label(s, "图表 / 图片 / 数据\n占位区域",
          cr, ct, cw, ch,
          size=Pt(16), color=MUTED)

# ════════════════════════════════════════════════════════════
# 4. 结尾页
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
cover_base(s)

tx = W * 0.38 + Inches(0.55)
tw = W * 0.60 - Inches(0.5)
label(s, "感谢聆听",
      tx, H * 0.18, tw, Inches(1.4),
      size=Pt(46), bold=True, color=DARK, align=PP_ALIGN.LEFT)
label(s, "THANK  YOU",
      tx, H * 0.18 + Inches(1.4), tw, Inches(0.85),
      size=Pt(28), color=CYAN, align=PP_ALIGN.LEFT)
rect(s, tx, H * 0.18 + Inches(2.32), Inches(2.2), Inches(0.06), fill=CYAN)
label(s, "如有问题，欢迎交流探讨",
      tx, H * 0.18 + Inches(2.48), tw, Inches(0.55),
      size=Pt(14), color=SIDEBAR, align=PP_ALIGN.LEFT)
label(s, "邮箱：example@university.edu.cn",
      tx, H * 0.76, tw, Inches(0.42),
      size=Pt(12), color=RGBColor(0x50, 0x78, 0xA0), align=PP_ALIGN.LEFT)

# ── 保存 ──────────────────────────────────────────────────
out = r"C:\Users\rdft1\OneDrive\Desktop\统计建模\nav_bar_template.pptx"
prs.save(out)
print("已保存：", out)
